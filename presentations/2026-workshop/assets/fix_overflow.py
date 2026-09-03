"""
Iteratively shrinks body/caption text on any slide whose content overflows
the safe content area. Overflow is detected by rendering the deck to PDF via
LibreOffice and inspecting text-block bounding boxes with PyMuPDF -- pandoc's
pptx writer doesn't auto-shrink text to fit a placeholder, so any slide with
enough content silently bleeds past the bottom edge. Titles, the
footer/page-number boxes, and the yellow callout shapes are left untouched.
Repeats until no slide overflows or a max-iteration cap is hit.
"""
import re
import subprocess
import sys
import tempfile
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

SOFFICE = r"C:\Program Files\LibreOffice\program\soffice.exe"
MAX_ITERATIONS = 6
MIN_FONT_PT = 12
FOOTER_MARGIN_PT = 40  # keep content this far above the slide bottom
CALLOUT_BUFFER_PT = 4  # keep content this far above a yellow callout box
EMU_PER_PT = 12700

# pandoc names code-block/inline-code runs "Courier" -- a font neither Windows
# nor macOS actually has installed under that exact name, so each renderer
# silently substitutes something different (LibreOffice's substitute happens
# to be narrower than PowerPoint's), making LibreOffice-measured overflow
# checks unreliable for code text. Normalizing to "Courier New" (a real font
# on both platforms) makes what we measure match what viewers actually see.
CODE_FONT_ALIASES = {"Courier"}
CODE_FONT_REPLACEMENT = "Courier New"

# Slide 9's code block still overflowed in real PowerPoint even after the
# Courier New fix -- LibreOffice's font metrics apparently still don't match
# PowerPoint's closely enough for this slide to trust the render-based loop
# alone. Cap it explicitly (~18% below its prior 16pt) as extra insurance.
# Expressed as a cap (not a relative cut) so re-running this script is safe:
# it only pulls oversized text down, never keeps shrinking on repeat runs.
EXTRA_SAFETY_CAPS_PT = {
    "Resulting project layout (under the hood)": 13,
}


def apply_extra_safety_caps(prs):
    for slide in prs.slides:
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0]
                break
        cap = EXTRA_SAFETY_CAPS_PT.get(title)
        if cap is None:
            continue
        for shape in slide.shapes:
            if is_protected(shape):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None or run.font.size.pt > cap:
                        run.font.size = Pt(cap)


def fix_code_font(prs):
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name in CODE_FONT_ALIASES:
                        run.font.name = CODE_FONT_REPLACEMENT
                        changed += 1
    return changed


def render_to_pdf(pptx_path, out_dir):
    def convert():
        subprocess.run(
            [SOFFICE, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
            check=True, capture_output=True,
        )

    pdf_path = out_dir / (Path(pptx_path).stem + ".pdf")
    convert()
    if not pdf_path.exists():
        convert()  # LibreOffice occasionally no-ops on first launch -- retry once
    return pdf_path


def is_footer_text(text):
    text = text.strip()
    return text.startswith("Copyright") or bool(re.fullmatch(r"\d+", text))


def callout_tops_by_slide(pptx_path):
    """Maps slide number -> top (pt) of its highest yellow callout box, if any."""
    prs = Presentation(pptx_path)
    tops = {}
    for i, slide in enumerate(prs.slides, 1):
        callout_tops = [
            shape.top / EMU_PER_PT
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        ]
        if callout_tops:
            tops[i] = min(callout_tops)
    return tops


def find_overflowing_slides(pdf_path, pptx_path):
    callout_tops = callout_tops_by_slide(pptx_path)
    doc = fitz.open(str(pdf_path))
    overflowing = set()
    for i, page in enumerate(doc, 1):
        safe_bottom = page.rect.height - FOOTER_MARGIN_PT
        if i in callout_tops:
            safe_bottom = min(safe_bottom, callout_tops[i] - CALLOUT_BUFFER_PT)
        for block in page.get_text("blocks"):
            x0, y0, x1, y1, text = block[:5]
            if is_footer_text(text):
                continue
            if i in callout_tops and y0 >= callout_tops[i] - CALLOUT_BUFFER_PT:
                continue  # this block is the callout's own text
            if y1 > safe_bottom:
                overflowing.add(i)
                break
    doc.close()
    return overflowing


def is_protected(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return True  # yellow callout boxes -- already sized deliberately
    if not shape.has_text_frame:
        return True
    text = shape.text_frame.text.strip()
    if not text or is_footer_text(text):
        return True
    if shape.is_placeholder and shape.placeholder_format.idx == 0:
        return True  # title
    return False


def shrink_slide(slide):
    for shape in slide.shapes:
        if is_protected(shape):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                cur = run.font.size
                new_pt = 22 if cur is None else max(cur.pt - 2, MIN_FONT_PT)
                run.font.size = Pt(new_pt)


def main():
    path = Path(sys.argv[1] if len(sys.argv) > 1 else "slides.pptx").resolve()

    prs = Presentation(path)
    changed = fix_code_font(prs)
    if changed:
        prs.save(path)
        print(f"Normalized {changed} run(s) from 'Courier' to '{CODE_FONT_REPLACEMENT}'")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        overflowing = set()
        converged = False
        for iteration in range(1, MAX_ITERATIONS + 1):
            pdf_path = render_to_pdf(path, tmp_dir)
            overflowing = find_overflowing_slides(pdf_path, path)
            if not overflowing:
                print(f"No overflow after {iteration - 1} shrink pass(es).")
                converged = True
                break
            print(f"Iteration {iteration}: overflow on slides {sorted(overflowing)}")
            prs = Presentation(path)
            for slide_num in overflowing:
                shrink_slide(prs.slides[slide_num - 1])
            prs.save(path)
        if not converged:
            print(f"Reached max iterations ({MAX_ITERATIONS}); remaining overflow: {sorted(overflowing)}")

    prs = Presentation(path)
    apply_extra_safety_caps(prs)
    prs.save(path)


if __name__ == "__main__":
    sys.exit(main())
