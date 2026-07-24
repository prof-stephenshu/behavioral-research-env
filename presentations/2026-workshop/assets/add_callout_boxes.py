"""
Adds yellow highlight callout boxes to two punchline slides ("The reveal"
and "Takeaways"), matching the reference deck's convention of a bold,
bordered yellow box for a key statement. Shrinks the existing text
placeholder slightly on each slide first to free clean space above the
footer, rather than overlapping existing content.
"""
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

YELLOW = RGBColor(0xFF, 0xFF, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)


def add_yellow_box(slide, left, top, width, height, text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = YELLOW
    box.line.color.rgb = BLACK
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BLACK
    return box


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else "slides.pptx"
    prs = Presentation(path)

    for i, slide in enumerate(prs.slides, 1):
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0]
                break

        if title == "The reveal":
            intro = find_shape(slide, "Text Placeholder 3")
            # Placeholders inheriting size from the layout report no explicit
            # width/height until one is set -- capture both dimensions BEFORE
            # changing either, and set all four (left/top/width/height)
            # together, or python-pptx will zero out the untouched dimension.
            orig_left, orig_top, orig_width = intro.left, intro.top, intro.width
            intro.left, intro.top, intro.width, intro.height = orig_left, orig_top, orig_width, Emu(2700000)
            add_yellow_box(
                slide,
                left=orig_left, top=Emu(3856326), width=orig_width, height=Emu(750000),
                text="We didn't design this to match a published study. It just did.",
            )
            print(f"Added callout to slide {i}: {title}")

        elif title == "Takeaways":
            body = find_shape(slide, "Content Placeholder 2")
            orig_left, orig_top, orig_width = body.left, body.top, body.width
            body.left, body.top, body.width, body.height = orig_left, orig_top, orig_width, Emu(2700000)
            add_yellow_box(
                slide,
                left=orig_left, top=Emu(3980151), width=orig_width, height=Emu(650000),
                text="Synthetic samples are a tool for iteration — not a substitute for real data.",
            )
            print(f"Added callout to slide {i}: {title}")

    prs.save(path)


if __name__ == "__main__":
    sys.exit(main())
