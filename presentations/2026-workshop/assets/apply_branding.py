"""
Post-processes slides.pptx to add a consistent footer (copyright, bottom-left)
and slide number (bold navy, bottom-right) to every slide, matching the
visual convention of the reference course deck (footer + numbered slides),
without reproducing any institution-specific logos or marks.
"""
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1F, 0x4E, 0x79)
GRAY = RGBColor(0x59, 0x59, 0x59)

COPYRIGHT_TEXT = "Copyright © 2026 Stephen Shu"


def add_footer(slide, slide_number, slide_width, slide_height):
    footer_box = slide.shapes.add_textbox(Inches(0.3), slide_height - Inches(0.4), Inches(6), Inches(0.3))
    tf = footer_box.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = COPYRIGHT_TEXT
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY

    number_box = slide.shapes.add_textbox(slide_width - Inches(0.9), slide_height - Inches(0.45), Inches(0.6), Inches(0.35))
    tf2 = number_box.text_frame
    tf2.margin_top = 0
    tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(slide_number)
    run2.font.size = Pt(12)
    run2.font.bold = True
    run2.font.color.rgb = NAVY


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else "slides.pptx"
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        add_footer(slide, i, prs.slide_width, prs.slide_height)
    prs.save(path)
    print(f"Applied footer/page-number branding to {len(prs.slides)} slides in {path}")


if __name__ == "__main__":
    sys.exit(main())
